import pptx
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR,MSO_AUTO_SIZE
from pptx.util import Cm,Pt
from pptx.dml.color import RGBColor

#https://blog.csdn.net/weixin_42750611/article/details/108029796#comments_15566677
prs = Presentation(r'F:\Python练习\办公自动化\积水框胶测试汇总.pptx')
for i,slide in enumerate(prs.slides):
    if i==1:
        for j,shape in enumerate(slide.shapes):
            print('*****************')
            print(shape)
            if shape.has_text_frame:    #判断是否有文字结构内容                
                #shape.text='hahahhahahahhaha'
                text_frame=shape.text_frame
                print(text_frame)
                for p in text_frame.paragraphs:
                    print(p)
                    p.text='6666666666我太厉害了hdeidjedjsahdiashjdkasjdlksajkldfhdjksfhklajdlkasjdklashfjksadf'
            #if shape.has_picture_frame:    #判断是否有文字结构内容
            #    print('chart')        
                    #add_picture
                    p.font.name = '黑体'
                    p.font.size= Pt(16)
                    p.font.color.rgb = RGBColor(0,0,0)
                text_frame.margin_bottom = Cm(0.1) #下边距
                text_frame.margin_left = 0 #左边距
                text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM # 对齐文本方式：底端对齐
                text_frame.word_wrap = True # 框中的文字自动换行
                #print(shape.text)

            else:
                print('none')
        left = top = width = height =Cm(3)
        text_box = slide.shapes.add_textbox(left,top,width,height)
        tf = text_box.text_frame
        tf.text = "这是一段文本框里的文字"
        tf.margin_bottom = Cm(0.1) #下边距
        tf.margin_left = 0 #左边距
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM # 对齐文本方式：底端对齐
        tf.word_wrap = True # 框中的文字自动换行

        left=Cm(0.45)
        top=Cm(6.52)
        #pic=slide.shapes.add_picture('总.png',left,top)
        height = Cm(5)
        pic =slide.shapes.add_picture('总.png',left,top, height=height,width=height)
        rows,cols = 4,2
        left = top =Cm(5)
        width = Cm(18)
        height = Cm(3)

        table = slide.shapes.add_table(rows,cols,left,top,width,height).table
# 可以修改列宽、行高
        table.columns[0].width = Cm(6)
        table.columns[1].width = Cm(4)
        table.rows[0].height =Cm(2)

        data = [
        ['姓名','成绩'],
        ['李雷',99],
        ['韩梅梅', 92],
        ['马东梅', 92],
        ]
        for row in range(rows):
            for col in range(cols):
                table.cell(row,col).text =str(data[row][col])
          
#.font.name 字体名称（可以直接设定为中文字体）
#.font.bold 是否加粗
#.font.italic 是否斜体
#.font.color 字体颜色
#.font.size 字体大小
prs.save('积水框胶测试汇总内容的写入.pptx')