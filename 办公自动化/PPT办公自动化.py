import pptx
#prs = Presentation()新建ppt对象
#Slide：幻灯片，就是演示文稿中每一页的页面。for slide in prs.slides
#Shape：方框，在每页幻灯片内插入的方框，可以是形状，也可以是文本框。for shape in slide.shapes
#Run：文字块，一般为较少字符。
#Paragraph：段落，通常有序号ㆍ、1.等

from pptx import Presentation
prs = Presentation(r'F:\Python练习\办公自动化\积水框胶测试汇总.pptx')
for slide in prs.slides:
    print(slide)
#prs = Presentation()
#prs.slides.add_slide(prs.slide_layouts[0])
print(enumerate(prs.slides))
for i,slide in enumerate(prs.slides):#enumerate() 函数用于将一个可遍历的数据对象(如列表、元组或字符串)组合为一个索引序列，同时列出数据和数据下标，一般用在 for 循环当中
    #print(i)

#print((enumerate(prs.slides))[1])
#for slide in prs.slides:
    print('*****************')
    for shape in slide.shapes:
        print(shape)
        if shape.has_text_frame:
            text_frame = shape.text_frame
            print(text_frame)
            print(text_frame.text)
            #for paragragh in text_frame.paragragh:
            #    print(paragragh)
            #    print(paragragh.text)
                
"""
注意：
该方法同样也直接获取Shpae中的文字内容；
但是这个更灵活，先获取每个Shape，然后在获取每个Shape中的paragraph；
我们可以针对paragraph，写一个判断条件，只获取第几个paragraph；
"""

prs.save('PPT办公自动化1.pptx')